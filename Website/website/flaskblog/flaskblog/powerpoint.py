from pptx import Presentation
from pptx.util import Inches
import os

class Powerpoint(object):


    ROOT_DIRECTORY = os.getcwd()

    def __init__(self, pptx, des_directory, save_name):
        self.pptx = pptx
        self.des_directory = des_directory
        self.save_name = save_name
    def create_powerpoint_slide(self):
        print('Creating powerpoint')
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Hello, World!"
        subtitle.text = "python-pptx was here!"

        prs.save(self.save_name)

    def open_powerpoint_slide(self, slide):
        prs = Presentation(slide)
        # use prs.save(name) to save the powerpoint file with the given name
        return prs

    def move_slide(self, presentation, old_index, new_index):
            xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
            print(xml_slides)
            slides = list(xml_slides)
            print(slides)
            xml_slides.remove(slides[old_index])
            xml_slides.insert(new_index, slides[old_index])
            #presentation.save('test.pptx')

    # identifier_array will be a list of strings of the path (including the name) of the png file for the identifier
    def add_images(self, slide_numbers, identifier_array):
        # initilaise pptx variables
        prs = self.open_powerpoint_slide(self.pptx)
        total_page_num = len(prs.slides)

        #if last page has ending => all good; else => insret manualy, shouldnt do this as we need the id_array to be updated
        if(not(total_page_num + 1 in slide_numbers)):   # we know that we want to insert the id after the last slide hence why total page num + 1
            print('Last page {} not in slide_numbers {} , now adding it'.format(total_page_num + 1, slide_numbers))
            slide_numbers.append(total_page_num+1)  # might have to do something special here
    
        # check that the positions array = identifier array length
        if (not len(slide_numbers) == len(identifier_array)):
            print("Slide_numbers array length = {}\nIdentifier Array Lenght = {}\nThese are NOT equal, thus returning".format(len(slide_numbers),len(identifier_array)))
            return

        # offset is needed because if we want to insert into the original at pos[1,2] then we want to put the identifiers
        # into [1,(iden_1),2,(iden_2),...]
        offset = 0
        # loop over PDF and insert identifier in slots
        for i in range(len(slide_numbers) - 1): # since cant add last slide no point in trying hence why (-1)
            page_position = slide_numbers[i]
            img =  identifier_array[i]
            print('Page Position: {}, Image directory: {}'.format(page_position,img))
            self.add_image(page_position + offset, img)
            offset += 1
        # note that add image will not append the last slide, thus
        last_slide_idx =  len(identifier_array) - 1
        self.append_image(identifier_array[last_slide_idx])
       
    def append_image(self,img):
        prs = self.open_powerpoint_slide(self.pptx)
        slide_height = prs.slide_height
        slide_width = prs.slide_width
        left = top = Inches(0)
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img, left, top, height=slide_height, width=slide_width)
        save_file = os.path.join(self.des_directory,self.pptx)
        prs.save(save_file)

    # slide number starts at 1 and will end at the length of the slides
    def add_image(self, slide_number, image_path):
        print('Adding image at slide {} using image path {}'.format(slide_number,image_path))
        img_path = image_path

        # might need to open up a pre-existing presentation
        prs = self.open_powerpoint_slide(self.pptx)

        slide_height = prs.slide_height
        slide_width = prs.slide_width
        left = top = Inches(0)
        print('Slide height: {}, slide width: {}'.format(slide_height,slide_width))
    
        slides = prs.slides
        total_slides = len(slides)
        print ('# of slides: {}'.format(total_slides))

        for slide in slides:
            slide_position = slides.index(slide)
            if (slide_number == (slide_position + 1)):
                print('Adding image to this slide number {}\nOld index:{}, New index: {}'.format((slides.index(slide)+1),total_slides,slide_position))
                # create a new slide
                blank_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(blank_slide_layout)
                # update total slides as weve just added one
                total_slides += 1
                # add the pic to the new slide
                pic = slide.shapes.add_picture(img_path, left, top, height=slide_height, width=slide_width)
                # move this slide to the correct position
                self.move_slide(prs,total_slides - 1 ,slide_position)    # total_slides -1 because we indent at 0
            
            #print('slide number {}'.format(slides.index(slide)+1))
    
        save_file = os.path.join(self.des_directory,self.pptx)
        prs.save(save_file)


    def save_pptx(self, save_name):
        prs = Presentation(self.pptx)
        save_file = os.path.join(self.des_directory,save_name)
        prs.save(save_file)

    def get_total_slides(self):
        prs = self.open_powerpoint_slide(self.pptx)
        total_page_num = len(prs.slides)
        return total_page_num

def test_add_image():
    pptx = 'test_images.pptx'
    powerpoint = Powerpoint(pptx)
    powerpoint.add_image(1,'Assassins creed.jpg',path,'test.pptx')

# works
def test_add_images_1():
    pptx = 'test_images.pptx'
    powerpoint = Powerpoint(pptx,os.getcwd(),pptx)
    slide_numbers = [1,3,4,8]
    insert_slide_positions = [x+1 for x in slide_numbers]
    id_array = ['Assassins creed.jpg', 'Assassins creed.jpg', 'Assassins creed.jpg', 'Assassins creed.jpg']
    des_directory = os.getcwd()
    
    powerpoint.add_images(insert_slide_positions, id_array)
    # expectation: slide 1 - txt, slide 2 - img, slide 3 - text, slide - 4 - text, slide 5 - img, slide 6 - text, slide 7- img, slide 8 - txt, slide 9 -txt, slide 10 - txt, slide 11 - txt, slide 12 - txt

# works
def test_add_images_2():
    pptx = 'test_images_2.pptx'
    powerpoint = Powerpoint(pptx,os.getcwd(),pptx)
    slide_numbers = [0,3,5,6,8]
    insert_slide_positions = [x+1 for x in slide_numbers]
    id_array = ['338105.jpg', '3801049-beach-sunset-wallpaper-desktop.jpg', '445384.jpg', 'Assassins creed.jpg','beaches-wallpaper-13.jpg']
    des_directory = os.getcwd()
    
    powerpoint.add_images(insert_slide_positions, id_array)
    # expectation: slide 1 - txt, slide 2 - txt, slide 3 - img, slide - 4 - text, slide 5 - img, slide 6 - text, slide 7- txt, slide 8 - img, slide 9 -txt, slide 10 - img, slide 11 - txt, slide 12 - txt, slide 13 - img


def main():
    test_add_images_2()

if __name__ == '__main__':
    main()