
'''
Workflow:
1) get pptx
2) magic
3) return zip file with segmented pptx files
4) ... apply water

... Do stuff ...
p.s. ilie is straight
and james is not 
...
'''
import os
import shutil
from pptx import Presentation


# delete slide id link from chain of slides and link to next one in pptx rID
def delete_slide(prs, slide):
    # Make dictionary with necessary information
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


# start slide included, end slide included
def segment_topic(prs, start_slide, end_slide,title=None):
	prs_temp  = prs
	slides_temp = prs_temp.slides
	title = 'from_{}_to_{}.pptx'.format(start_slide,end_slide) if title==None else '{}from_{}_to_{}.pptx'.format(title[0],title[1],title[2])
	print('len(slides) = {}\nstart_slide={}\nend_slide={}'.format(len(slides_temp),start_slide,end_slide))
	for slide in slides_temp:
		slide_nr = slides_temp.index(slide)
		print('slide_nr = {}'.format(slide_nr))
		# delete slides if not desired 
		if(start_slide > slide_nr or slide_nr > end_slide):
			print('\nDELETING SLIDE\nstart_slide = {} slide_nr = {} end_slide = {}'.format(start_slide,slide_nr,end_slide))
			delete_slide(prs_temp,slide)
			if start_slide!=0:
				end_slide -= 1
				start_slide -= 1
		else:
			# save the slide
			print('\nSAVING SLIDE\nstart_slide = {} slide_nr = {} end_slide = {}'.format(start_slide,slide_nr,end_slide))
			pass
	print('saving file')
	if len(prs_temp.slides) > 1:
		prs_temp.save(title)	
	else:
		print('empty file')
	print('-'*40)


'''
Returns: zip with pptx files
Takes: pptx source directory
	   pptx destination directory
	   lsit where to segment
'''
def segment_pptx_from_list(source_file, destination='del_dir\\', topic_endings=[]):
	start_slide = 0
	if not os.path.exists(destination):
		os.makedirs(destination)
	for idx in topic_endings:
		prs = Presentation(source_file)
		end_slide = idx
		title = [destination,start_slide,end_slide]
		segment_topic(prs,start_slide,end_slide,title)
		print('\nSAVING SLIDE\nstart_slide = {} end_slide = {}'.format(start_slide,end_slide))
		start_slide = end_slide
	# in case you want the last slides and not included
	prs = Presentation('randomFile.pptx')
	title = [destination,end_slide, len(prs.slides)]
	segment_topic(prs,0,len(prs.slides),title)
	#in case you want he whole presentation as well
	prs = Presentation('randomFile.pptx')
	title = [destination,0, len(prs.slides)]
	segment_topic(prs,0,len(prs.slides),title)
	# make the zip file
	zip_file_name='pptx_topics.zip'
	shutil.make_archive('pptx_topics', 'zip', destination)
	# remove zipped directory
	shutil.rmtree(destination)
	#return zip file location
	zip_path = os.path.dirname(os.path.realpath(zip_file_name)) + '\\' + zip_file_name
	return zip_path


'''
This is an example of how the workflow is to be used
'''
def main():
	topic_endings = [3, 5, 7, 11, 17]
	zip_path = segment_pptx_from_list('randomFile.pptx','to_zip\\',topic_endings)
	print('zip_path={}'.format(zip_path))


if __name__=='__main__':
	main()
